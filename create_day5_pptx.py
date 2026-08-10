import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.abspath(__file__))
DASH = os.path.join(BASE, "dashboard")
OUT  = os.path.join(BASE, "reports", "Day5_Advanced_Analytics.pptx")

BG_DARK = RGBColor(0x0d, 0x11, 0x17)
PANEL   = RGBColor(0x16, 0x1b, 0x22)
BLUE    = RGBColor(0x4f, 0xc3, 0xf7)
GREEN   = RGBColor(0x56, 0xd3, 0x64)
AMBER   = RGBColor(0xe3, 0xb3, 0x41)
WHITE   = RGBColor(0xe6, 0xed, 0xf3)
GRAY    = RGBColor(0x8b, 0x94, 0x9e)
PURPLE  = RGBColor(0xbc, 0x8c, 0xff)
RED     = RGBColor(0xf8, 0x51, 0x49)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

def add_slide():
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_DARK
    return s

def txb(slide, text, x, y, w, h, size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def rect(slide, x, y, w, h, fill=PANEL, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_img(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)

def header(slide, title, subtitle=""):
    rect(slide, 0, 0, W, Inches(0.75), fill=PANEL)
    txb(slide, title, Inches(0.3), Inches(0.08), Inches(10), Inches(0.45),
        size=22, bold=True, color=BLUE)
    if subtitle:
        txb(slide, subtitle, Inches(0.3), Inches(0.50), Inches(12), Inches(0.25),
            size=10, color=GRAY)

# ── SLIDE 1: Title ────────────────────────────────────────────────────────────
s = add_slide()
rect(s, Inches(1.5), Inches(1.9), Inches(10.3), Inches(0.06), fill=BLUE)
txb(s, "Mutual Fund Analysis - Day 5",
    Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.8),
    size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "Advanced Analytics + Risk Metrics",
    Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6),
    size=24, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
txb(s, "VaR/CVaR  |  Rolling Sharpe  |  Cohort Analysis  |  SIP Continuity  |  Recommender  |  HHI",
    Inches(0.5), Inches(2.1), Inches(12.3), Inches(0.4),
    size=12, color=GRAY, align=PP_ALIGN.CENTER)

tasks = [
    ("Task 1", "VaR & CVaR", BLUE),
    ("Task 2", "Rolling Sharpe", GREEN),
    ("Task 3", "Cohort Analysis", AMBER),
    ("Task 4", "SIP Continuity", PURPLE),
    ("Task 5", "Recommender", RED),
    ("Task 6", "Sector HHI", BLUE),
]
for i, (tag, desc, col) in enumerate(tasks):
    bx = Inches(0.4 + i * 2.1)
    rect(s, bx, Inches(2.9), Inches(1.9), Inches(0.85), fill=PANEL, line=col)
    txb(s, tag,  bx, Inches(2.93), Inches(1.9), Inches(0.3), size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s, desc, bx, Inches(3.25), Inches(1.9), Inches(0.3), size=9,  color=WHITE, align=PP_ALIGN.CENTER)

txb(s, "GitHub: github.com/JASWANTH1726/mf-analysis  |  Commit: 26d14ec",
    Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.3),
    size=10, color=GRAY, align=PP_ALIGN.CENTER)
txb(s, "Tools: Python | Pandas | NumPy | Matplotlib | Seaborn | Jupyter",
    Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.3),
    size=10, color=GRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 2: Tasks Overview ───────────────────────────────────────────────────
s = add_slide()
header(s, "Day 5 - Tasks Overview", "6 Advanced Analytics Tasks + 5 Insights + 6 Deliverables")

items = [
    (BLUE,   "Task 1 - Historical VaR (95%) & CVaR",
             "5th percentile of daily return distribution. CVaR = mean of returns below VaR. Computed for all 40 schemes. Saved to var_cvar_report.csv"),
    (GREEN,  "Task 2 - Rolling 90-Day Sharpe Ratio",
             "returns.rolling(90).mean() / returns.rolling(90).std() x sqrt(252). Plotted for top 5 scorecard funds. Saved as rolling_sharpe_chart.png"),
    (AMBER,  "Task 3 - Investor Cohort Analysis",
             "Group by first SIP year. Compute avg SIP amount, total invested, top fund preference per cohort. Saved to cohort_analysis.csv"),
    (PURPLE, "Task 4 - SIP Continuity Analysis",
             "For folios with 3+ SIPs, compute avg gap between dates. Flag folios with gap > 35 days as at-risk. Saved to sip_continuity.csv"),
    (RED,    "Task 5 - Fund Recommender",
             "Input: risk appetite (Low/Moderate/High). Output: top 3 funds by Sharpe within matching risk_grade. Saved as recommender.py"),
    (BLUE,   "Task 6 - Sector HHI Concentration",
             "HHI = sum(weight_i^2) per fund. High HHI = concentrated portfolio. Compared across all equity funds. Saved to hhi_concentration.csv"),
]

for i, (col, title, desc) in enumerate(items):
    y = Inches(0.85 + i * 1.08)
    rect(s, Inches(0.3), y, Inches(12.7), Inches(0.95), fill=PANEL, line=col)
    txb(s, title, Inches(0.5), y + Inches(0.05), Inches(12.3), Inches(0.32),
        size=11, bold=True, color=col)
    txb(s, desc,  Inches(0.5), y + Inches(0.40), Inches(12.3), Inches(0.40),
        size=9, color=WHITE)

# ── SLIDE 3: VaR & CVaR ───────────────────────────────────────────────────────
s = add_slide()
header(s, "Task 1 - Historical VaR (95%) & CVaR", "All 40 Schemes  |  Daily Return Distribution  |  Tail Risk")

rect(s, Inches(0.3), Inches(0.85), Inches(12.7), Inches(1.5), fill=PANEL, line=BLUE)
txb(s, "Formula", Inches(0.5), Inches(0.90), Inches(12.3), Inches(0.3),
    size=11, bold=True, color=BLUE)
txb(s, "VaR 95% = np.percentile(daily_returns, 5)     |     CVaR 95% = mean(returns where return <= VaR)",
    Inches(0.5), Inches(1.22), Inches(12.3), Inches(0.3), size=10, color=WHITE)
txb(s, "Interpretation: On a bad day (5th percentile), VaR tells max expected loss. CVaR = average loss beyond that threshold.",
    Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.3), size=9, color=GRAY)

results = [
    ("Highest VaR (Most Risky)",  "Nippon India Flexi Cap Fund",      "VaR: -0.3843%", "CVaR: -0.55%", RED),
    ("Lowest VaR (Safest)",       "Aditya Birla Flexi Cap Fund",      "VaR: -0.1392%", "CVaR: -0.20%", GREEN),
    ("High Risk Grade Avg",       "Very High / High risk funds",      "VaR: ~-0.30%",  "CVaR: ~-0.45%", AMBER),
    ("Low Risk Grade Avg",        "Low / Moderate risk funds",        "VaR: ~-0.15%",  "CVaR: ~-0.22%", BLUE),
]
for i, (label, fund, var, cvar, col) in enumerate(results):
    col_x = Inches(0.3 + (i % 2) * 6.5)
    row_y = Inches(2.55 + (i // 2) * 2.2)
    rect(s, col_x, row_y, Inches(6.2), Inches(1.9), fill=PANEL, line=col)
    txb(s, label, col_x + Inches(0.15), row_y + Inches(0.08), Inches(5.9), Inches(0.3),
        size=10, bold=True, color=col)
    txb(s, fund,  col_x + Inches(0.15), row_y + Inches(0.42), Inches(5.9), Inches(0.3),
        size=10, color=WHITE)
    txb(s, var + "     " + cvar,
        col_x + Inches(0.15), row_y + Inches(0.78), Inches(5.9), Inches(0.3),
        size=11, bold=True, color=col)
    txb(s, "Saved: var_cvar_report.csv (40 rows)",
        col_x + Inches(0.15), row_y + Inches(1.15), Inches(5.9), Inches(0.3),
        size=8, color=GRAY)

# ── SLIDE 4: Rolling Sharpe Chart ─────────────────────────────────────────────
s = add_slide()
header(s, "Task 2 - Rolling 90-Day Sharpe Ratio", "Top 5 Scorecard Funds  |  Annualised  |  2015-2024")

img_path = os.path.join(DASH, "rolling_sharpe_chart.png")
add_img(s, img_path, Inches(0.3), Inches(0.85), Inches(12.7), Inches(5.5))

txb(s, "Formula: (rolling_mean(90) - RF) / rolling_std(90) x sqrt(252)     RF = 6.5% / 252 daily",
    Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.3),
    size=9, color=GRAY, align=PP_ALIGN.CENTER)
txb(s, "Quantum Balanced Advantage is the only fund crossing positive Sharpe in 2023-24",
    Inches(0.3), Inches(6.85), Inches(12.7), Inches(0.3),
    size=9, color=GREEN, align=PP_ALIGN.CENTER)

# ── SLIDE 5: Cohort Analysis ──────────────────────────────────────────────────
s = add_slide()
header(s, "Task 3 - Investor Cohort Analysis", "Grouped by First SIP Year  |  Avg SIP  |  Total Invested  |  Top Fund")

cohort_data = [
    ("2020", "156", "Rs 14,268", "Rs 23.97L", "DSP Flexi Cap Fund"),
    ("2021", "140", "Rs 14,109", "Rs 21.45L", "ICICI Short Duration"),
    ("2022", "162", "Rs 11,874", "Rs 20.31L", "ICICI Short Duration"),
    ("2023", "157", "Rs 16,948", "Rs 27.46L", "Nippon Flexi Cap Fund"),
    ("2024", "145", "Rs 12,517", "Rs 18.15L", "Quantum Balanced Adv"),
]

headers_row = ["Cohort Year", "Investors", "Avg SIP Amt", "Total Invested", "Top Fund"]
col_widths   = [1.5, 1.5, 2.0, 2.0, 4.5]
col_x        = [0.3, 1.8, 3.3, 5.3, 7.3]

rect(s, Inches(0.3), Inches(0.9), Inches(11.5), Inches(0.45), fill=BLUE)
for j, (hdr, cx) in enumerate(zip(headers_row, col_x)):
    txb(s, hdr, Inches(cx + 0.05), Inches(0.93), Inches(col_widths[j]), Inches(0.35),
        size=10, bold=True, color=BG_DARK)

for i, row in enumerate(cohort_data):
    y = Inches(1.45 + i * 0.7)
    fill = PANEL if i % 2 == 0 else RGBColor(0x1c, 0x21, 0x28)
    rect(s, Inches(0.3), y, Inches(11.5), Inches(0.6), fill=fill)
    for j, (val, cx) in enumerate(zip(row, col_x)):
        col_color = AMBER if j == 3 else WHITE
        txb(s, val, Inches(cx + 0.05), y + Inches(0.12), Inches(col_widths[j]), Inches(0.35),
            size=10, color=col_color)

rect(s, Inches(0.3), Inches(5.1), Inches(12.7), Inches(0.6), fill=PANEL, line=AMBER)
txb(s, "Key Insight: 2023 cohort invests the most (Rs 27.46L total, avg Rs 16,948/SIP). Newer investors commit larger amounts.",
    Inches(0.5), Inches(5.18), Inches(12.3), Inches(0.35), size=10, color=AMBER)

# ── SLIDE 6: SIP Continuity ───────────────────────────────────────────────────
s = add_slide()
header(s, "Task 4 - SIP Continuity Analysis", "Folio-level Gap Analysis  |  At-Risk Flag  |  Gap > 35 Days")

metrics = [
    ("Method",          "Sort SIP transactions by folio_no + date. Compute gap between consecutive SIPs.", BLUE),
    ("Active Folios",   "Folios with 3+ SIP transactions selected for analysis", GREEN),
    ("At-Risk Flag",    "avg_gap_days > 35 days = at-risk investor (skipping monthly SIPs)", RED),
    ("Business Use",    "At-risk investors should receive nudge campaigns to resume SIPs before lapse", AMBER),
    ("Output File",     "sip_continuity.csv - columns: folio_no, avg_gap_days, at_risk (True/False)", PURPLE),
]

for i, (label, desc, col) in enumerate(metrics):
    y = Inches(0.9 + i * 1.2)
    rect(s, Inches(0.3), y, Inches(12.7), Inches(1.0), fill=PANEL, line=col)
    txb(s, label, Inches(0.5), y + Inches(0.05), Inches(3.0), Inches(0.35),
        size=11, bold=True, color=col)
    txb(s, desc,  Inches(3.6), y + Inches(0.25), Inches(9.2), Inches(0.45),
        size=10, color=WHITE)

# ── SLIDE 7: Fund Recommender ─────────────────────────────────────────────────
s = add_slide()
header(s, "Task 5 - Fund Recommender", "Input: Risk Appetite  |  Output: Top 3 Funds by Sharpe  |  recommender.py")

rect(s, Inches(0.3), Inches(0.85), Inches(12.7), Inches(0.5), fill=PANEL, line=BLUE)
txb(s, "Usage:  python recommender.py     OR     from recommender import recommend; recommend('Moderate')",
    Inches(0.5), Inches(0.93), Inches(12.3), Inches(0.3), size=10, color=WHITE)

recs = [
    ("Low Risk",      GREEN,  [
        ("Edelweiss Flexi Cap Fund",        "Low",      "-0.648", "49.0"),
        ("Kotak Mahindra Liquid Fund",      "Low",      "-0.774", "62.4"),
        ("Tata Small Cap Fund",             "Low",      "-0.887", "37.2"),
    ]),
    ("Moderate Risk", AMBER,  [
        ("Sundaram Liquid Fund",            "Moderate", "-0.541", "53.8"),
        ("Aditya Birla Flexi Cap Fund",     "Moderate", "-0.543", "40.9"),
        ("DSP Flexi Cap Fund",              "Moderate", "-0.591", "47.9"),
    ]),
    ("High Risk",     RED,    [
        ("Quantum Balanced Advantage",      "Very High","+0.002", "61.4"),
        ("PGIM India Mid Cap Fund",         "High",     "-0.410", "54.1"),
        ("Sundaram Liquid Fund",            "Moderate", "-0.541", "53.8"),
    ]),
]

for i, (appetite, col, funds) in enumerate(recs):
    bx = Inches(0.3 + i * 4.35)
    rect(s, bx, Inches(1.5), Inches(4.1), Inches(0.45), fill=col)
    txb(s, appetite, bx, Inches(1.53), Inches(4.1), Inches(0.35),
        size=12, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    for j, (fname, rg, sharpe, score) in enumerate(funds):
        fy = Inches(2.05 + j * 1.6)
        rect(s, bx, fy, Inches(4.1), Inches(1.45), fill=PANEL, line=col)
        txb(s, fname[:30],  bx + Inches(0.1), fy + Inches(0.05), Inches(3.9), Inches(0.35),
            size=10, bold=True, color=col)
        txb(s, "Risk Grade: " + rg, bx + Inches(0.1), fy + Inches(0.42), Inches(3.9), Inches(0.25),
            size=9, color=GRAY)
        txb(s, "Sharpe: " + sharpe + "   Score: " + score,
            bx + Inches(0.1), fy + Inches(0.72), Inches(3.9), Inches(0.3),
            size=10, bold=True, color=WHITE)

# ── SLIDE 8: Sector HHI ───────────────────────────────────────────────────────
s = add_slide()
header(s, "Task 6 - Sector HHI Concentration", "Herfindahl-Hirschman Index  |  Equity Funds  |  Portfolio Concentration")

rect(s, Inches(0.3), Inches(0.85), Inches(12.7), Inches(0.5), fill=PANEL, line=BLUE)
txb(s, "Formula: HHI = sum(weight_i ^ 2)   |   High HHI (>0.15) = Concentrated   |   Low HHI (<0.08) = Diversified",
    Inches(0.5), Inches(0.93), Inches(12.3), Inches(0.3), size=10, color=WHITE)

hhi_data = [
    ("Kotak Mahindra Mid Cap Fund",  "0.2192", "High",     "Banking", RED),
    ("HDFC Large Cap Fund",          "0.1885", "High",     "Auto",    RED),
    ("HDFC ELSS Fund",               "0.1864", "High",     "FMCG",    RED),
    ("Kotak Mahindra Liquid Fund",   "0.1841", "High",     "Energy",  RED),
    ("Edelweiss Flexi Cap Fund",     "0.1435", "Moderate", "Banking", AMBER),
    ("Tata Short Duration Fund",     "0.1405", "Moderate", "IT",      AMBER),
    ("Tata Small Cap Fund",          "0.1393", "Moderate", "Banking", AMBER),
    ("Edelweiss Balanced Advantage", "0.1239", "Moderate", "Energy",  AMBER),
]

hdrs = ["Fund Name", "HHI Score", "Concentration", "Top Sector"]
cw   = [5.5, 1.8, 2.2, 2.2]
cx   = [0.3, 5.8, 7.6, 9.8]

rect(s, Inches(0.3), Inches(1.45), Inches(11.9), Inches(0.4), fill=BLUE)
for j, (hdr, x) in enumerate(zip(hdrs, cx)):
    txb(s, hdr, Inches(x + 0.05), Inches(1.48), Inches(cw[j]), Inches(0.3),
        size=10, bold=True, color=BG_DARK)

for i, (fname, hhi, conc, sector, col) in enumerate(hhi_data):
    y = Inches(1.95 + i * 0.62)
    fill = PANEL if i % 2 == 0 else RGBColor(0x1c, 0x21, 0x28)
    rect(s, Inches(0.3), y, Inches(11.9), Inches(0.55), fill=fill)
    for j, (val, x) in enumerate(zip([fname, hhi, conc, sector], cx)):
        c = col if j in [1, 2] else WHITE
        txb(s, val, Inches(x + 0.05), y + Inches(0.1), Inches(cw[j]), Inches(0.3),
            size=9, color=c)

# ── SLIDE 9: 5 Advanced Insights ─────────────────────────────────────────────
s = add_slide()
header(s, "5 Advanced Insights", "From Advanced_Analytics.ipynb")

insights = [
    (BLUE,   "Insight 1 - Highest VaR Risk",
             "Nippon India Flexi Cap has highest VaR (-0.38%). Aditya Birla Flexi Cap is safest (-0.14%). High risk grade funds consistently show VaR below -0.30%."),
    (GREEN,  "Insight 2 - Investor Cohorts",
             "2023 cohort invests the most (Rs 27.46L, avg Rs 16,948/SIP). Newer investors commit larger amounts. Fund preference shifted from DSP (2020) to ICICI (2021-22) to Nippon (2023)."),
    (AMBER,  "Insight 3 - SIP Continuity",
             "Folios with avg gap > 35 days are at-risk of lapsing. These investors need nudge campaigns. High continuity rate (>80%) means disciplined investor base."),
    (PURPLE, "Insight 4 - Rolling Sharpe & Market Cycles",
             "All top 5 funds had negative Sharpe during 2022 market correction. Quantum Balanced Advantage is the only fund crossing positive Sharpe in 2023-24 due to defensive allocation."),
    (RED,    "Insight 5 - Sector Concentration Risk",
             "Kotak Mid Cap has highest HHI (0.22) - Banking concentrated. HDFC Large Cap & ELSS also High HHI. Investors seeking diversification should prefer HHI < 0.08 funds."),
]

for i, (col, title, desc) in enumerate(insights):
    row, ci = divmod(i, 2)
    if i == 4:
        x = Inches(0.3)
        w = Inches(12.7)
    else:
        x = Inches(0.3 + ci * 6.5)
        w = Inches(6.2)
    y = Inches(0.85 + row * 2.2)
    rect(s, x, y, w, Inches(1.95), fill=PANEL, line=col)
    txb(s, title, x + Inches(0.15), y + Inches(0.08), w - Inches(0.3), Inches(0.35),
        size=11, bold=True, color=col)
    txb(s, desc,  x + Inches(0.15), y + Inches(0.5),  w - Inches(0.3), Inches(1.2),
        size=9, color=WHITE)

# ── SLIDE 10: Deliverables & Submission ───────────────────────────────────────
s = add_slide()
header(s, "Deliverables & Submission Details", "Day 5 - Advanced Analytics + Risk Metrics")

rect(s, Inches(0.3), Inches(0.85), Inches(12.7), Inches(6.4), fill=PANEL)

files = [
    ("advanced_analytics.py",                  "Main script - all 6 tasks",                    BLUE),
    ("recommender.py",                          "Standalone fund recommender",                  BLUE),
    ("var_cvar_report.csv",                     "VaR & CVaR for all 40 schemes",                GREEN),
    ("cohort_analysis.csv",                     "Investor cohort metrics by first SIP year",    GREEN),
    ("sip_continuity.csv",                      "Folio-level SIP gap + at-risk flag",           GREEN),
    ("hhi_concentration.csv",                   "Sector HHI for all equity funds",              GREEN),
    ("notebooks/Advanced_Analytics.ipynb",      "Full notebook with 5 markdown insights",       AMBER),
    ("dashboard/rolling_sharpe_chart.png",      "Rolling 90-day Sharpe chart - top 5 funds",   AMBER),
]

for i, (fname, desc, col) in enumerate(files):
    row, ci = divmod(i, 2)
    x = Inches(0.5 + ci * 6.4)
    y = Inches(1.0 + row * 1.3)
    rect(s, x, y, Inches(6.1), Inches(1.1), fill=RGBColor(0x1c, 0x21, 0x28), line=col)
    txb(s, fname, x + Inches(0.12), y + Inches(0.06), Inches(5.8), Inches(0.35),
        size=10, bold=True, color=col)
    txb(s, desc,  x + Inches(0.12), y + Inches(0.45), Inches(5.8), Inches(0.4),
        size=9, color=WHITE)

txb(s, "GitHub: github.com/JASWANTH1726/mf-analysis  |  Commit: 26d14ec  |  Branch: main",
    Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.3),
    size=10, color=GRAY, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides))
