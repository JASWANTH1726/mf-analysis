import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE  = os.path.dirname(__file__)
DASH  = os.path.join(BASE, "dashboard")
OUT   = os.path.join(BASE, "reports", "Day4_Dashboard.pptx")

# ── colours ───────────────────────────────────────────────────────────────────
BG_DARK  = RGBColor(0x0d, 0x11, 0x17)
PANEL    = RGBColor(0x16, 0x1b, 0x22)
BLUE     = RGBColor(0x4f, 0xc3, 0xf7)
GREEN    = RGBColor(0x56, 0xd3, 0x64)
AMBER    = RGBColor(0xe3, 0xb3, 0x41)
WHITE    = RGBColor(0xe6, 0xed, 0xf3)
GRAY     = RGBColor(0x8b, 0x94, 0x9e)
PURPLE   = RGBColor(0xbc, 0x8c, 0xff)
RED      = RGBColor(0xf8, 0x51, 0x49)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # completely blank

def add_slide():
    s = prs.slides.add_slide(blank)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_DARK
    return s

def txb(slide, text, x, y, w, h, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def rect(slide, x, y, w, h, fill=PANEL, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_img(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)

def header(slide, title, subtitle=""):
    rect(slide, 0, 0, W, Inches(0.75), fill=PANEL)
    txb(slide, title, Inches(0.3), Inches(0.08), Inches(10), Inches(0.45),
        size=24, bold=True, color=BLUE, align=PP_ALIGN.LEFT)
    if subtitle:
        txb(slide, subtitle, Inches(0.3), Inches(0.52), Inches(12), Inches(0.25),
            size=11, color=GRAY, align=PP_ALIGN.LEFT)

def kpi(slide, x, y, w, h, label, value, color=BLUE):
    rect(slide, x, y, w, h, fill=PANEL, line=color)
    txb(slide, value, x, y + Inches(0.08), w, Inches(0.35),
        size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    txb(slide, label, x, y + Inches(0.42), w, Inches(0.22),
        size=9, color=GRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, W, H, fill=BG_DARK)
rect(s, Inches(1.5), Inches(1.8), Inches(10.3), Inches(0.06), fill=BLUE)

txb(s, "Mutual Fund Analysis — Day 4",
    Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.8),
    size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txb(s, "Dashboard Development",
    Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6),
    size=26, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

txb(s, "4-Page Interactive Dashboard  |  PNG + PDF Export  |  40 Schemes  |  Dark Theme",
    Inches(0.5), Inches(2.1), Inches(12.3), Inches(0.4),
    size=13, color=GRAY, align=PP_ALIGN.CENTER)

# deliverable boxes
items = [
    ("Page 1", "Industry Overview", BLUE),
    ("Page 2", "Fund Performance", GREEN),
    ("Page 3", "Investor Analytics", AMBER),
    ("Page 4", "SIP & Market Trends", PURPLE),
]
for i, (tag, desc, col) in enumerate(items):
    bx = Inches(0.5 + i * 3.1)
    rect(s, bx, Inches(2.8), Inches(2.8), Inches(0.9), fill=PANEL, line=col)
    txb(s, tag,  bx, Inches(2.85), Inches(2.8), Inches(0.35), size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s, desc, bx, Inches(3.2),  Inches(2.8), Inches(0.35), size=10, color=WHITE, align=PP_ALIGN.CENTER)

txb(s, "GitHub: github.com/JASWANTH1726/mf-analysis  |  Commit: fb9b695",
    Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.35),
    size=11, color=GRAY, align=PP_ALIGN.CENTER)

txb(s, "Tools: Python  |  Matplotlib  |  Seaborn  |  Pandas  |  PIL  |  SQLite",
    Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.35),
    size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Agenda / What Was Built
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
header(s, "Day 4 — What Was Built", "Dashboard Development  |  4 Pages  |  PNG + PDF")

tasks = [
    ("create_dashboard.py", "Single Python script generating all 4 dashboard pages + PDF", BLUE),
    ("Page 1 — Industry Overview",
     "KPI cards (AUM, SIP, Folios, Schemes) + AUM trend + AUM by AMC + Category pie + Risk donut", GREEN),
    ("Page 2 — Fund Performance",
     "Return vs Risk scatter (bubble=AUM, colour=score) + Scorecard table top 10 + NAV vs benchmark + Alpha bar", AMBER),
    ("Page 3 — Investor Analytics",
     "Txn type donut + KYC bar + State transaction bar + Monthly volume grouped bar + Avg amount by fund house", PURPLE),
    ("Page 4 — SIP & Market Trends",
     "SIP vs Nifty 50 dual-axis + Category inflow heatmap + Net inflow by category + SIP YoY dual-axis", RED),
    ("Dashboard.pdf", "All 4 pages combined into a single PDF using PIL + matplotlib PdfPages", BLUE),
]

for i, (title, desc, col) in enumerate(tasks):
    y = Inches(0.9 + i * 1.0)
    rect(s, Inches(0.3), y, Inches(12.7), Inches(0.85), fill=PANEL, line=col)
    txb(s, f"✦  {title}", Inches(0.5), y + Inches(0.05), Inches(12.3), Inches(0.3),
        size=12, bold=True, color=col)
    txb(s, desc, Inches(0.7), y + Inches(0.38), Inches(12.0), Inches(0.35),
        size=10, color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Page 1: Industry Overview
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
header(s, "Page 1 — Industry Overview", "KPI Cards  |  AUM Trend  |  Fund House  |  Category Mix  |  Risk Grade")

img_path = os.path.join(DASH, "page1_industry_overview.png")
add_img(s, img_path, Inches(0.3), Inches(0.85), Inches(12.7), Inches(6.0))

txb(s, "40 Schemes  |  20 Fund Houses  |  2022–2024 AUM Data  |  Dark Theme (#0d1117)",
    Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.35),
    size=9, color=GRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Page 2: Fund Performance
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
header(s, "Page 2 — Fund Performance Analytics", "Return vs Risk  |  Scorecard Top 10  |  NAV vs Benchmark  |  Alpha Bar")

img_path = os.path.join(DASH, "page2_fund_performance.png")
add_img(s, img_path, Inches(0.3), Inches(0.85), Inches(12.7), Inches(6.0))

txb(s, "Scorecard #1: Mirae Short Duration (70.06)  |  Top Alpha: Quantum Balanced Advantage (44.30%)",
    Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.35),
    size=9, color=GRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Page 3: Investor Analytics
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
header(s, "Page 3 — Investor Analytics", "Transactions  |  KYC Status  |  State Distribution  |  Monthly Volume")

img_path = os.path.join(DASH, "page3_investor_analytics.png")
add_img(s, img_path, Inches(0.3), Inches(0.85), Inches(12.7), Inches(6.0))

txb(s, "1,998 Transactions  |  SIP / Lumpsum / Redemption  |  KYC Verified / Pending / Rejected",
    Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.35),
    size=9, color=GRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Page 4: SIP & Market Trends
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
header(s, "Page 4 — SIP & Market Trends", "SIP vs Nifty 50  |  Category Heatmap  |  Net Inflows  |  YoY Trend")

img_path = os.path.join(DASH, "page4_sip_market_trends.png")
add_img(s, img_path, Inches(0.3), Inches(0.85), Inches(12.7), Inches(6.0))

txb(s, "SIP Inflows vs Nifty 50 dual-axis  |  Category heatmap by year  |  Net inflow by sub-category",
    Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.35),
    size=9, color=GRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Key Insights
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
header(s, "Key Insights from Dashboard", "Fund Scorecard  |  Alpha/Beta  |  AUM  |  Investor Behaviour")

insights = [
    (BLUE,   "Scorecard #1",   "Mirae Short Duration (118989) — Score 70.06 — balanced across CAGR, Sharpe, Alpha, TER"),
    (GREEN,  "Top Alpha",      "Quantum Balanced Advantage (121020) — Alpha 44.30% — only fund with positive Sharpe (+0.002)"),
    (AMBER,  "CAGR Leader",    "PGIM Mid Cap (121009) — CAGR 3Y 29.59% — highest 3-year compounded return"),
    (PURPLE, "AUM Insight",    "Top 10 fund houses hold ~80% of total AUM; HDFC, SBI, ICICI dominate"),
    (RED,    "Investor Mix",   "SIP transactions dominate (~60%); KYC Verified accounts drive 85%+ of inflows"),
    (BLUE,   "Market Trend",   "SIP inflows show strong YoY growth; Nifty 50 correlation visible in monthly data"),
    (GREEN,  "Category Mix",   "Equity funds lead (40%); Debt & Hybrid each ~30%; Balanced Advantage growing"),
    (AMBER,  "Risk Profile",   "Moderate risk grade dominates (50%+); Very High risk funds show highest alpha"),
]

for i, (col, tag, text) in enumerate(insights):
    row, col_idx = divmod(i, 2)
    x = Inches(0.3 + col_idx * 6.5)
    y = Inches(0.9 + row * 1.55)
    rect(s, x, y, Inches(6.2), Inches(1.3), fill=PANEL, line=col)
    txb(s, tag,  x + Inches(0.15), y + Inches(0.08), Inches(5.9), Inches(0.35),
        size=12, bold=True, color=col)
    txb(s, text, x + Inches(0.15), y + Inches(0.45), Inches(5.9), Inches(0.65),
        size=10, color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Technical Architecture
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
header(s, "Technical Architecture", "Data Pipeline  |  Analytics  |  Visualisation  |  Export")

layers = [
    ("Data Layer",        "12 CSVs → data_cleaning.py → 10 clean CSVs → db_load.py → bluestock_mf.db (SQLite star schema)", BLUE),
    ("Analytics Layer",   "performance_analytics.py → daily_returns, CAGR, Sharpe, Sortino, Alpha/Beta (OLS), Max Drawdown, Scorecard", GREEN),
    ("Dashboard Layer",   "create_dashboard.py → matplotlib (Agg backend) + seaborn → 4 × PNG (150 dpi, 20×11.25 in)", AMBER),
    ("Export Layer",      "PIL + PdfPages → Dashboard.pdf (4 pages combined)  |  Git push → github.com/JASWANTH1726/mf-analysis", PURPLE),
]

for i, (layer, desc, col) in enumerate(layers):
    y = Inches(0.9 + i * 1.5)
    rect(s, Inches(0.3), y, Inches(2.2), Inches(1.2), fill=col)
    txb(s, layer, Inches(0.3), y + Inches(0.35), Inches(2.2), Inches(0.5),
        size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    rect(s, Inches(2.7), y, Inches(10.3), Inches(1.2), fill=PANEL, line=col)
    txb(s, desc, Inches(2.9), y + Inches(0.3), Inches(10.0), Inches(0.6),
        size=11, color=WHITE)
    if i < 3:
        txb(s, "▼", Inches(1.1), y + Inches(1.25), Inches(0.6), Inches(0.25),
            size=14, color=col, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Deliverables & File Structure
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
header(s, "Deliverables & File Structure", "All outputs committed to GitHub")

files = [
    ("create_dashboard.py",                    "Main dashboard script — 4 pages + PDF", BLUE),
    ("dashboard/page1_industry_overview.png",  "Industry KPIs, AUM trend, category mix", GREEN),
    ("dashboard/page2_fund_performance.png",   "Return/risk scatter, scorecard, alpha bar", GREEN),
    ("dashboard/page3_investor_analytics.png", "Txn donut, KYC, state, monthly volume", GREEN),
    ("dashboard/page4_sip_market_trends.png",  "SIP vs Nifty, heatmap, net inflows, YoY", GREEN),
    ("dashboard/Dashboard.pdf",                "All 4 pages combined into single PDF", AMBER),
    ("fund_scorecard.csv",                     "40 funds scored — top: Mirae Short Duration (70.06)", PURPLE),
    ("alpha_beta.csv",                         "OLS alpha/beta for 40 funds — top alpha: Quantum (44.30%)", PURPLE),
    ("bluestock_mf.db",                        "SQLite DB — dim_fund(40), fact_nav(139,760), fact_txn(1,998)", BLUE),
]

for i, (fname, desc, col) in enumerate(files):
    row, ci = divmod(i, 2)
    x = Inches(0.3 + ci * 6.5)
    y = Inches(0.9 + row * 1.3)
    rect(s, x, y, Inches(6.2), Inches(1.1), fill=PANEL, line=col)
    txb(s, fname, x + Inches(0.12), y + Inches(0.06), Inches(5.9), Inches(0.35),
        size=10, bold=True, color=col)
    txb(s, desc,  x + Inches(0.12), y + Inches(0.45), Inches(5.9), Inches(0.45),
        size=9, color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Submission Details
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
header(s, "Submission Details", "Day 4 — Dashboard Development")

rect(s, Inches(0.3), Inches(0.85), Inches(12.7), Inches(6.4), fill=PANEL)

details = [
    ("Project",       "Mutual Fund Analysis — Bluestock Fintech Internship"),
    ("Day",           "Day 4 — Dashboard Development"),
    ("Task",          "Build a 4-page matplotlib dashboard replacing Power BI"),
    ("GitHub Repo",   "https://github.com/JASWANTH1726/mf-analysis"),
    ("Latest Commit", "fb9b695 — Dashboard: 4-page PNG + PDF complete"),
    ("Branch",        "main"),
    ("Files Added",   "create_dashboard.py, 4× page PNGs, Dashboard.pdf"),
    ("Dataset",       "40 schemes, 139,760 NAV rows, 1,998 transactions, 1,440 AUM records"),
    ("DB",            "bluestock_mf.db — SQLite star schema (6 tables)"),
    ("Top Fund",      "Mirae Short Duration (118989) — Scorecard Score: 70.06"),
    ("Top Alpha",     "Quantum Balanced Advantage (121020) — Alpha: 44.30%"),
    ("Tech Stack",    "Python 3.x | Pandas | Matplotlib | Seaborn | PIL | SQLite | Git"),
]

for i, (key, val) in enumerate(details):
    row, ci = divmod(i, 2)
    x = Inches(0.5 + ci * 6.4)
    y = Inches(1.0 + row * 0.85)
    txb(s, f"{key}:", x, y, Inches(1.6), Inches(0.35),
        size=10, bold=True, color=BLUE)
    txb(s, val, x + Inches(1.65), y, Inches(4.6), Inches(0.35),
        size=10, color=WHITE)

txb(s, "Submitted by: Jaswanth  |  Bluestock Fintech Internship  |  Day 4",
    Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.3),
    size=10, color=GRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
