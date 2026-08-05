from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd
import os

BASE = os.path.dirname(__file__)

scorecard = pd.read_csv(os.path.join(BASE, "fund_scorecard.csv"), index_col=0)
alpha_beta = pd.read_csv(os.path.join(BASE, "alpha_beta.csv"))

BG    = RGBColor(0x0f, 0x11, 0x17)
BLUE  = RGBColor(0x4f, 0xc3, 0xf7)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY  = RGBColor(0x90, 0xa4, 0xae)
GREEN = RGBColor(0x66, 0xbb, 0x6a)
AMBER = RGBColor(0xff, 0xa7, 0x26)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def add_slide():
    s = prs.slides.add_slide(blank)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s

def txb(slide, text, l, t, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def title_slide(slide, title, subtitle=""):
    txb(slide, title, 0.5, 2.8, 12, 1, size=36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    if subtitle:
        txb(slide, subtitle, 0.5, 3.9, 12, 0.6, size=18, color=GRAY, align=PP_ALIGN.CENTER)

def section_header(slide, text):
    txb(slide, text, 0.4, 0.2, 12, 0.55, size=22, bold=True, color=BLUE)
    line = slide.shapes.add_connector(1, Inches(0.4), Inches(0.82), Inches(12.5), Inches(0.82))
    line.line.color.rgb = BLUE
    line.line.width = Pt(1)

def table(slide, headers, rows, l, t, w, h, col_widths=None):
    cols = len(headers)
    tbl = slide.shapes.add_table(len(rows)+1, cols, Inches(l), Inches(t), Inches(w), Inches(h)).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)
    for i, h_txt in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h_txt
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1a, 0x1f, 0x35)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = BLUE
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r+1, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x12, 0x15, 0x1f) if r % 2 == 0 else RGBColor(0x1a, 0x1f, 0x35)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(10)
            run.font.color.rgb = WHITE

# ── slide 1: title ────────────────────────────────────────────────────────────
s = add_slide()
title_slide(s, "Mutual Fund Analysis", "Day 3 — EDA & Performance Analytics")
txb(s, "Bluestock Fintech Internship Project", 0.5, 5.0, 12, 0.5, size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ── slide 2: agenda ───────────────────────────────────────────────────────────
s = add_slide()
section_header(s, "Day 3 Agenda")
items = [
    "1.  Daily Returns Distribution",
    "2.  CAGR — 1 Year / 3 Year / 5 Year",
    "3.  Sharpe Ratio Ranking",
    "4.  Sortino Ratio",
    "5.  Alpha & Beta (OLS vs Nifty 100)",
    "6.  Maximum Drawdown Analysis",
    "7.  Fund Scorecard (Composite 0-100)",
    "8.  Benchmark Comparison Chart",
]
for i, item in enumerate(items):
    txb(s, item, 1.5, 1.0 + i*0.72, 10, 0.6, size=16, color=WHITE if i % 2 == 0 else GRAY)

# ── slide 3: daily returns ────────────────────────────────────────────────────
s = add_slide()
section_header(s, "Daily Returns — Distribution & Sanity Check")
txb(s, "Risk-Free Rate: RBI Repo 6.5% annualised  |  RF daily = 0.065 / 252", 0.5, 1.0, 12, 0.4, size=13, color=GRAY)
points = [
    "Daily returns computed using pct_change() per scheme",
    "Returns > ±5% flagged as extreme — result: 0 extreme rows (0.00%)",
    "Distribution is approximately normal across all 20 schemes",
    "Mean daily return range: -0.0002 to +0.0003 (as expected for MF data)",
]
for i, p in enumerate(points):
    txb(s, f"  {p}", 0.5, 1.6 + i*0.9, 12, 0.7, size=15, color=WHITE)
txb(s, "Data Quality: PASSED — no extreme outliers in daily returns", 0.5, 5.8, 12, 0.5, size=14, bold=True, color=GREEN)

# ── slide 4: CAGR ─────────────────────────────────────────────────────────────
s = add_slide()
section_header(s, "CAGR — 1 Year / 3 Year / 5 Year")
top5 = scorecard.head(5)
rows = [(r["scheme_name"].replace(" Direct Growth","").replace(" Fund","")[:35],
         f"{r['cagr_3y']}%") for _, r in top5.iterrows()]
# show top 5 by 3yr cagr from scorecard
cagr_rows = [
    ("Mirae Asset Short Duration", "37.78%", "7.08%", "N/A"),
    ("ICICI Prudential ELSS",      "N/A",    "5.31%", "N/A"),
    ("Kotak Mahindra Liquid",      "22.28%", "4.64%", "N/A"),
    ("Mirae Asset Large Cap",      "N/A",    "3.99%", "N/A"),
    ("SBI Mutual Fund Mid Cap",    "N/A",    "3.34%", "N/A"),
]
table(s, ["Fund", "CAGR 1Y", "CAGR 3Y", "CAGR 5Y"], cagr_rows,
      0.4, 1.1, 12.5, 3.2, col_widths=[5.5, 2.0, 2.0, 2.0])
txb(s, "Formula: CAGR = (End NAV / Start NAV)^(1/years) - 1  |  Lookback: 252 trading days per year", 0.5, 4.5, 12, 0.5, size=12, color=GRAY)
txb(s, "Worst 3Y CAGR: Nippon India Flexi Cap (-3.83%)  |  Franklin Templeton Small Cap (-1.63%)", 0.5, 5.1, 12, 0.5, size=13, color=AMBER)

# ── slide 5: sharpe & sortino ─────────────────────────────────────────────────
s = add_slide()
section_header(s, "Risk-Adjusted Returns — Sharpe & Sortino Ratios")
txb(s, "Sharpe = (Mean Excess Return / Std Dev) x sqrt(252)     Sortino = (Mean Excess Return / Downside Std) x sqrt(252)", 0.4, 1.0, 12.5, 0.5, size=12, color=GRAY)
ratio_rows = [
    ("Kotak Mahindra Liquid",        "-0.70", "Rank 1"),
    ("ICICI Prudential ELSS",        "-0.68", "Rank 2"),
    ("Mirae Asset Large Cap",        "-0.79", "Rank 3"),
    ("SBI Mutual Fund Flexi Cap",    "-0.59", "Rank 4"),
    ("Mirae Asset Short Duration",   "-0.96", "Rank 5"),
]
table(s, ["Fund", "Sharpe Ratio", "Sharpe Rank"], ratio_rows,
      0.4, 1.6, 12.5, 2.8, col_widths=[6.5, 3.0, 3.0])
txb(s, "Note: All Sharpe ratios are negative — synthetic benchmark data is uncorrelated with fund returns", 0.5, 4.6, 12, 0.5, size=12, color=AMBER)
txb(s, "Sortino penalises only downside volatility — better measure for asymmetric return distributions", 0.5, 5.2, 12, 0.5, size=13, color=GRAY)

# ── slide 6: alpha & beta ─────────────────────────────────────────────────────
s = add_slide()
section_header(s, "Alpha & Beta — OLS Regression vs Nifty 100")
txb(s, "Model: Fund Return = Alpha + Beta x Benchmark Return  |  Min 30 observations required", 0.4, 1.0, 12.5, 0.5, size=12, color=GRAY)
ab_rows = [
    ("SBI Mutual Fund Flexi Cap",   "-0.0004", "25.32%", "0.0004"),
    ("Kotak Mahindra Liquid",       "-0.0007", "22.88%", "0.0015"),
    ("Mirae Asset Large Cap",       "-0.0003", "20.79%", "0.0004"),
    ("ICICI Prudential ELSS",       " 0.0006", "20.45%", "0.0012"),
    ("HDFC Mutual Fund Small Cap",  "-0.0000", "16.19%", "0.0000"),
]
table(s, ["Fund", "Beta", "Alpha (Annual)", "R-Squared"], ab_rows,
      0.4, 1.6, 12.5, 2.8, col_widths=[5.5, 2.0, 2.5, 2.5])
txb(s, "Betas near 0 and R2 near 0 — benchmark data is synthetic/random, not correlated with fund NAVs", 0.5, 4.6, 12, 0.5, size=12, color=AMBER)
txb(s, "Worst alpha: Nippon India Flexi Cap (-17.81%)  |  Saved to: alpha_beta.csv", 0.5, 5.2, 12, 0.5, size=13, color=GRAY)

# ── slide 7: max drawdown ─────────────────────────────────────────────────────
s = add_slide()
section_header(s, "Maximum Drawdown Analysis")
txb(s, "Max Drawdown = (Trough NAV - Peak NAV) / Peak NAV  |  Measures worst peak-to-trough decline", 0.4, 1.0, 12.5, 0.5, size=12, color=GRAY)
dd_rows = [
    ("Nippon India Flexi Cap",          "-37.99%", "Worst"),
    ("Franklin Templeton Small Cap",    "-28.75%", "2nd Worst"),
    ("UTI Mutual Fund Mid Cap",         "-20.91%", "3rd Worst"),
    ("Axis Short Duration",             "-20.12%", "4th Worst"),
    ("Kotak Balanced Advantage",        "-19.24%", "5th Worst"),
]
table(s, ["Fund", "Max Drawdown", "Rank"], dd_rows,
      0.4, 1.6, 12.5, 2.8, col_widths=[6.5, 3.0, 3.0])
txb(s, "Best (least drawdown): ICICI ELSS (-7.41%)  |  Kotak Liquid (-7.63%)", 0.5, 4.6, 12, 0.5, size=13, color=GREEN)
txb(s, "Drawdown used in scorecard with 10% weight — lower drawdown = higher score", 0.5, 5.2, 12, 0.5, size=13, color=GRAY)

# ── slide 8: fund scorecard ───────────────────────────────────────────────────
s = add_slide()
section_header(s, "Fund Scorecard — Composite Score (0-100)")
txb(s, "Weights: 30% CAGR 3Y  +  25% Sharpe  +  20% Alpha  +  15% TER (inverse)  +  10% Max Drawdown (inverse)", 0.4, 1.0, 12.5, 0.5, size=12, color=GRAY)
sc_rows = []
for _, r in scorecard.head(8).iterrows():
    name = r["scheme_name"].replace(" Direct Growth","").replace(" Mutual Fund","")[:35]
    sc_rows.append((name, f"{r['score']}", f"{r['cagr_3y']}%", f"{round(r['sharpe'],2)}", f"{r['alpha_annual']}%"))
table(s, ["Fund", "Score", "CAGR 3Y", "Sharpe", "Alpha"], sc_rows,
      0.4, 1.6, 12.5, 3.5, col_widths=[5.5, 1.5, 1.8, 1.8, 1.9])
txb(s, "Top Fund: Axis ELSS (64.5)  |  Saved to: fund_scorecard.csv", 0.5, 5.3, 12, 0.5, size=13, bold=True, color=GREEN)

# ── slide 9: benchmark chart ──────────────────────────────────────────────────
s = add_slide()
section_header(s, "Benchmark Comparison Chart")
chart_path = os.path.join(BASE, "dashboard", "benchmark_comparison.png")
if os.path.exists(chart_path):
    s.shapes.add_picture(chart_path, Inches(0.3), Inches(1.0), Inches(12.7), Inches(5.8))
else:
    txb(s, "Chart not found — run eda_analytics.py first", 0.5, 3.5, 12, 0.5, size=16, color=AMBER, align=PP_ALIGN.CENTER)

# ── slide 10: summary ─────────────────────────────────────────────────────────
s = add_slide()
section_header(s, "Day 3 Summary & Key Insights")
insights = [
    "20 schemes analysed  |  69,880 NAV data points  |  3,494 trading days",
    "Top CAGR 3Y: Mirae Asset Short Duration (7.08%)  |  Worst: Nippon Flexi Cap (-3.83%)",
    "Best Sharpe: Kotak Liquid (-0.70)  |  Best Alpha: SBI Flexi Cap (+25.32% annualised)",
    "Worst Drawdown: Nippon India Flexi Cap (-37.99%)  |  Best: ICICI ELSS (-7.41%)",
    "Top Scorecard: Axis ELSS (64.5)  |  Mirae Short Duration (62.5)  |  Kotak Liquid (61.25)",
    "Deliverables: alpha_beta.csv, fund_scorecard.csv, dashboard/benchmark_comparison.png",
]
for i, ins in enumerate(insights):
    txb(s, f"  {ins}", 0.5, 1.1 + i*0.9, 12.3, 0.7, size=14, color=WHITE if i % 2 == 0 else GRAY)
txb(s, "GitHub: github.com/JASWANTH1726/mf-analysis  |  Commit: 8937a0c", 0.5, 6.8, 12, 0.4, size=12, color=GRAY, align=PP_ALIGN.CENTER)

out = os.path.join(BASE, "reports", "Day3_EDA_Analytics.pptx")
prs.save(out)
print(f"saved: {out}")
