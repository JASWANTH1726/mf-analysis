"""Generate the final Bluestock MF presentation deck."""
from pptx import Presentation
from pptx.util import Inches, Pt
import os

OUT = "Bluestock_MF_Presentation.pptx"

SLIDES = [
    ("Bluestock MF Capstone", "Final project presentation for mutual fund analytics and dashboard reporting."),
    ("Problem & Objective", "Problem: Identify top mutual fund schemes in India using data-driven analytics. Objective: Build a repeatable ETL and reporting pipeline that supports investment comparison and dashboard visualization."),
    ("Data Sources", "Data sources include scheme metadata, NAV history, benchmark series, SIP transactions, AUM, expense ratio, fund managers, and risk metrics. Live NAV is fetched from MFAPI."),
    ("Architecture", "The solution uses a structured ETL pipeline: dataset generation, data cleaning, database load, EDA, analytics, and deliverable generation with automated scripts."),
    ("EDA Highlights", "NAV distribution and volatility patterns across large-cap schemes. Relationship between expense ratio and risk-adjusted return."),
    ("EDA Highlights", "Correlation and portfolio composition analysis reveal sector concentration and risk profiles across schemes."),
    ("Performance Metrics", "CAGR, Sharpe ratio, Sortino ratio, alpha/beta, drawdown, and fund scorecard metrics show scheme performance across horizons."),
    ("Performance Metrics", "VaR/CVaR and SIP continuity analysis support downside risk evaluation and investor behavior insights."),
    ("Dashboard Screenshots", "Visual dashboards illustrate comparative fund performance, investor analytics, and portfolio risk snapshots."),
    ("Dashboard Screenshots", "The interactive reporting layout supports stakeholder review and executive decisions."),
    ("Key Findings", "Focus on risk-adjusted returns, expense management, sector concentration, and consistent SIP investor engagement for portfolio selection."),
    ("Thank You", "Thank you for reviewing the Bluestock MF capstone analysis. Questions welcome.")
]


def find_images():
    images = []
    for root, _, files in os.walk("dashboard"):
        for file in sorted(files):
            if file.lower().endswith((".png", ".jpg", ".jpeg")):
                images.append(os.path.join(root, file))
    return images[:4]


def add_text_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.text = content
    for paragraph in body.paragraphs:
        paragraph.font.size = Pt(18)


def add_image_slide(prs, title, image_path):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))


def build():
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Bluestock MF Capstone"
    slide.placeholders[1].text = "Final project presentation for mutual fund analytics and dashboard reporting."

    for title, content in SLIDES[1:]:
        add_text_slide(prs, title, content)

    images = find_images()
    for i, image_path in enumerate(images, start=1):
        add_image_slide(prs, f"Dashboard Screenshot {i}", image_path)

    while len(prs.slides) < 12:
        add_text_slide(prs, "Supporting Insight", "This slide supports the final analytical findings and conclusions.")

    prs.save(OUT)
    print(f"Generated presentation: {OUT}")


if __name__ == "__main__":
    build()
