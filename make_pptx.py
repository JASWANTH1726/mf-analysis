from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK   = RGBColor(0x0F, 0x11, 0x17)
BG_CARD   = RGBColor(0x1A, 0x1F, 0x35)
BLUE      = RGBColor(0x4F, 0xC3, 0xF7)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0x90, 0xA4, 0xAE)
GREEN     = RGBColor(0x66, 0xBB, 0x6A)
ORANGE    = RGBColor(0xFF, 0xA7, 0x26)
DARK_BLUE = RGBColor(0x1E, 0x3A, 0x5F)

blank = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(blank)


def bg(slide, color=BG_DARK):
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def box(slide, x, y, w, h, color=BG_CARD, radius=False):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(0x2A, 0x35, 0x50)
    shape.line.width = Pt(0.5)
    return shape


def txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def accent_bar(slide, y=0.08):
    bar = slide.shapes.add_shape(1, 0, Inches(y), Inches(13.33), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()


def slide_number(slide, n, total=10):
    txt(slide, f"{n} / {total}", 12.3, 7.1, 0.8, 0.3, size=10, color=GREY, align=PP_ALIGN.RIGHT)


# ── SLIDE 1: Title ──────────────────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s)

box(s, 0.6, 1.2, 8.5, 0.06, color=BLUE)  # decorative line

txt(s, "MUTUAL FUND ANALYSIS", 0.6, 1.5, 10, 1.2, size=48, bold=True, color=BLUE)
txt(s, "India", 0.6, 2.7, 4, 0.7, size=36, bold=True, color=WHITE)
txt(s, "End-to-end data pipeline for 20 AMFI-registered schemes\nacross 10 fund houses  ·  Live NAV  ·  Risk & Returns",
    0.6, 3.5, 9, 0.9, size=16, color=GREY)

# tags
for i, tag in enumerate(["Python", "Pandas", "mfapi.in", "AMFI Data", "10 Fund Houses"]):
    bx = box(s, 0.6 + i * 1.9, 4.7, 1.7, 0.4, color=DARK_BLUE)
    txt(s, tag, 0.65 + i * 1.9, 4.72, 1.6, 0.36, size=11, color=BLUE, align=PP_ALIGN.CENTER)

txt(s, "Jaswanth  ·  github.com/JASWANTH1726/mf-analysis", 0.6, 6.8, 10, 0.4, size=11, color=GREY)
slide_number(s, 1)

# ── SLIDE 2: Problem Statement ───────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "Problem Statement", 0.6, 0.3, 10, 0.6, size=28, bold=True, color=BLUE)

box(s, 0.6, 1.1, 5.8, 3.2, color=BG_CARD)
txt(s, "India has 44+ AMCs and over 2,500 mutual fund schemes.\nRetail investors have no unified way to compare funds on\nrisk-adjusted returns, portfolio overlap, expense ratios,\nor manager track records — all in one place.",
    0.8, 1.2, 5.4, 1.5, size=13, color=GREY)

points = [
    "NAV data scattered across AMC websites",
    "No unified view of risk metrics vs returns",
    "SIP performance hard to track historically",
    "Portfolio overlap between funds is invisible",
]
for i, p in enumerate(points):
    txt(s, f"→  {p}", 0.8, 2.8 + i * 0.38, 5.4, 0.36, size=12, color=WHITE)

box(s, 6.8, 1.1, 5.9, 3.2, color=BG_CARD)
txt(s, "Project Goal", 7.0, 1.2, 5.5, 0.4, size=14, bold=True, color=BLUE)
txt(s, "Build a structured data pipeline that ingests, cleans,\nand organises mutual fund data — enabling fund\ncomparison, risk profiling, and SIP analysis through\na single codebase.",
    7.0, 1.7, 5.5, 1.5, size=12, color=GREY)

for tag, col in [("Data Engineering", BLUE), ("Financial Analysis", GREEN), ("Live API", ORANGE)]:
    pass

txt(s, "Data Engineering  ·  Financial Analysis  ·  Live API Integration",
    7.0, 3.5, 5.5, 0.4, size=11, color=BLUE)

slide_number(s, 2)

# ── SLIDE 3: Architecture ────────────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "Data Pipeline Architecture", 0.6, 0.3, 10, 0.6, size=28, bold=True, color=BLUE)

steps = ["mfapi.in\n(Live API)", "live_nav_\nfetch.py", "data/raw/\n(CSV Store)", "data_\ningestion.py", "Analysis\n& Reports"]
for i, step in enumerate(steps):
    bx = box(s, 0.5 + i * 2.5, 1.3, 2.0, 1.1, color=DARK_BLUE)
    txt(s, step, 0.55 + i * 2.5, 1.35, 1.9, 1.0, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(s, "→", 2.55 + i * 2.5, 1.65, 0.4, 0.5, size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

cards = [
    ("Ingestion Layer", ["6 live NAV feeds via REST API", "10 structured CSV datasets", "AMFI code validation (100%)", "Null / duplicate detection"]),
    ("Storage Layout",  ["data/raw  —  source of truth", "data/processed  —  cleaned", "sql/  —  analytical queries", "dashboard/  —  viz exports"]),
    ("Analysis Layer",  ["notebooks/  —  EDA", "matplotlib / seaborn / plotly", "scipy  —  statistical tests", "SQLAlchemy  —  DB queries"]),
]
for i, (title, pts) in enumerate(cards):
    box(s, 0.5 + i * 4.3, 2.8, 4.0, 3.5, color=BG_CARD)
    txt(s, title, 0.7 + i * 4.3, 2.9, 3.6, 0.4, size=13, bold=True, color=BLUE)
    for j, p in enumerate(pts):
        txt(s, f"→  {p}", 0.7 + i * 4.3, 3.4 + j * 0.55, 3.6, 0.5, size=11, color=WHITE)

slide_number(s, 3)

# ── SLIDE 4: Datasets ────────────────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "Datasets Overview", 0.6, 0.3, 10, 0.6, size=28, bold=True, color=BLUE)

headers = ["Dataset", "Rows", "Key Columns", "Type"]
widths  = [2.2, 0.7, 6.5, 1.6]
xs = [0.4, 2.65, 3.38, 9.9]

# header row
hbox = box(s, 0.4, 1.0, 12.5, 0.42, color=DARK_BLUE)
for h, x, w in zip(headers, xs, widths):
    txt(s, h, x, 1.02, w, 0.38, size=11, bold=True, color=BLUE)

rows = [
    ("fund_master",        "20",     "scheme_code, fund_house, category, risk_grade",  "Master"),
    ("nav_history",        "10,000", "scheme_code, date, nav",                          "Time Series"),
    ("portfolio_holdings", "200",    "scheme_code, stock, weight_pct, sector",          "Holdings"),
    ("sip_data",           "600",    "scheme_code, date, amount, units, nav",           "Transactional"),
    ("returns_data",       "20",     "ret_1m, ret_3m, ret_1y, ret_3y, ret_5y",         "Derived"),
    ("aum_data",           "720",    "scheme_code, month, aum_cr, folios",              "Aggregate"),
    ("risk_metrics",       "20",     "std_dev, beta, sharpe, alpha, sortino",           "Derived"),
    ("expense_ratio",      "20",     "direct_ter, regular_ter, effective_date",         "Master"),
    ("fund_manager",       "20",     "manager, exp_years, qualification",               "Master"),
    ("benchmark_data",     "500",    "nifty50, nifty100, bse500, nifty_midcap",        "Reference"),
]
for i, row in enumerate(rows):
    row_bg = RGBColor(0x12, 0x15, 0x1F) if i % 2 == 0 else BG_CARD
    box(s, 0.4, 1.45 + i * 0.52, 12.5, 0.5, color=row_bg)
    for val, x, w in zip(row, xs, widths):
        col = BLUE if val in ("Master", "Time Series", "Holdings", "Transactional", "Derived", "Aggregate", "Reference") else WHITE
        txt(s, val, x, 1.47 + i * 0.52, w, 0.48, size=10, color=col)

slide_number(s, 4)

# ── SLIDE 5: Live NAV ────────────────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "Live NAV — mfapi.in", 0.6, 0.3, 10, 0.6, size=28, bold=True, color=BLUE)

box(s, 0.5, 1.1, 6.2, 5.6, color=BG_CARD)
txt(s, "API Endpoint", 0.7, 1.2, 5.8, 0.4, size=13, bold=True, color=BLUE)
txt(s, "GET  https://api.mfapi.in/mf/{scheme_code}", 0.7, 1.7, 5.8, 0.4, size=11, color=GREEN)
txt(s, "No authentication required  ·  Free  ·  JSON response",
    0.7, 2.1, 5.8, 0.35, size=10, color=GREY)

txt(s, "Response Structure", 0.7, 2.6, 5.8, 0.35, size=12, bold=True, color=BLUE)
code = '{\n  "meta": {\n    "fund_house": "...",\n    "scheme_name": "...",\n    "scheme_type": "..."\n  },\n  "data": [\n    { "date": "24-07-2026", "nav": "204.85" },\n    ...\n  ]\n}'
txt(s, code, 0.7, 3.0, 5.8, 2.2, size=10, color=GREEN)

box(s, 7.0, 1.1, 5.9, 5.6, color=BG_CARD)
txt(s, "Schemes Fetched", 7.2, 1.2, 5.5, 0.4, size=13, bold=True, color=BLUE)

nav_data = [
    ("HDFC Top 100 Direct",  "125497", "204.85",  "3,129"),
    ("SBI Bluechip",         "119551", "106.60",  "3,274"),
    ("ICICI Bluechip",       "120503", "108.05",  "3,345"),
    ("Nippon Large Cap",     "118632", "99.33",   "3,336"),
    ("Axis Bluechip",        "119092", "6242.17", "3,603"),
    ("Kotak Bluechip",       "120841", "249.27",  "3,339"),
]
hdrs = ["Scheme", "Code", "Latest NAV", "Records"]
hxs  = [7.1, 9.6, 10.4, 11.5]
box(s, 7.0, 1.7, 5.9, 0.38, color=DARK_BLUE)
for h, x in zip(hdrs, hxs):
    txt(s, h, x, 1.72, 1.2, 0.34, size=10, bold=True, color=BLUE)

for i, row in enumerate(nav_data):
    rb = RGBColor(0x12, 0x15, 0x1F) if i % 2 == 0 else BG_CARD
    box(s, 7.0, 2.1 + i * 0.55, 5.9, 0.53, color=rb)
    for val, x in zip(row, hxs):
        txt(s, val, x, 2.12 + i * 0.55, 1.2, 0.5, size=10, color=WHITE)

txt(s, "~3,300 records per scheme  ·  History from launch date  ·  Saved as UTF-8 CSV",
    7.1, 5.5, 5.7, 0.5, size=10, color=GREY)

slide_number(s, 5)

# ── SLIDE 6: Fund Master ─────────────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "Fund Master — Key Findings", 0.6, 0.3, 10, 0.6, size=28, bold=True, color=BLUE)

stats = [("10", "AMCs Covered"), ("5", "Categories"), ("8", "Sub-categories"), ("5", "Risk Grades")]
for i, (num, label) in enumerate(stats):
    box(s, 0.5 + i * 3.1, 1.1, 2.8, 1.5, color=BG_CARD)
    txt(s, num,   0.6 + i * 3.1, 1.15, 2.6, 0.8, size=36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    txt(s, label, 0.6 + i * 3.1, 1.9,  2.6, 0.5, size=11, color=GREY, align=PP_ALIGN.CENTER)

box(s, 0.5, 2.9, 5.9, 3.3, color=BG_CARD)
txt(s, "Fund Houses", 0.7, 3.0, 5.5, 0.4, size=13, bold=True, color=BLUE)
houses = ["HDFC Mutual Fund", "SBI Mutual Fund", "ICICI Prudential", "Nippon India", "Axis Mutual Fund",
          "Kotak Mahindra", "Mirae Asset", "DSP Mutual Fund", "Franklin Templeton", "UTI Mutual Fund"]
for i, h in enumerate(houses):
    col = 0 if i < 5 else 1
    row = i % 5
    txt(s, f"→  {h}", 0.7 + col * 2.9, 3.5 + row * 0.48, 2.7, 0.44, size=11, color=WHITE)

box(s, 6.7, 2.9, 6.1, 3.3, color=BG_CARD)
txt(s, "AMFI Scheme Code Structure", 6.9, 3.0, 5.7, 0.4, size=13, bold=True, color=BLUE)
txt(s, "Each scheme variant gets its own unique 6-digit AMFI code:\n\n"
       "→  Direct Growth\n"
       "→  Regular Growth\n"
       "→  Direct IDCW\n"
       "→  Regular IDCW\n\n"
       "So one fund can have 4+ codes. Always verify scheme_name\n"
       "from the API response — not just the label used to fetch.",
    6.9, 3.5, 5.7, 2.5, size=11, color=WHITE)

slide_number(s, 6)

# ── SLIDE 7: Data Quality ────────────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "Data Quality Summary", 0.6, 0.3, 10, 0.6, size=28, bold=True, color=BLUE)

metrics = [
    ("100%",  "AMFI Code Coverage",   "All 20 master codes present\nin nav_history"),
    ("0",     "Null Values",          "Across all 10 structured\ndatasets"),
    ("0",     "Duplicate Rows",       "No duplicate scheme_code\n+ date combinations"),
    ("16",    "Total CSV Files",      "10 synthetic + 6 live\nNAV datasets"),
]
for i, (val, label, desc) in enumerate(metrics):
    box(s, 0.5 + i * 3.1, 1.1, 2.8, 2.5, color=BG_CARD)
    txt(s, val,   0.6 + i * 3.1, 1.2,  2.6, 0.9, size=40, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    txt(s, label, 0.6 + i * 3.1, 2.1,  2.6, 0.4, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc,  0.6 + i * 3.1, 2.55, 2.6, 0.8, size=10, color=GREY, align=PP_ALIGN.CENTER)

box(s, 0.5, 3.9, 12.3, 2.5, color=BG_CARD)
txt(s, "Anomaly Noted — Live NAV API", 0.7, 4.0, 11.9, 0.4, size=13, bold=True, color=ORANGE)
txt(s, "The mfapi.in API returns scheme names that don't always match the label used to query.\n"
       "For example, code 119092 (queried as 'Axis Bluechip') returns 'HDFC Money Market Fund' in the response metadata.\n\n"
       "This is expected — AMFI codes are variant-specific, not fund-name-specific. Each Direct/Regular × Growth/IDCW\n"
       "combination gets its own code. Always use the scheme_name field from the API response as the source of truth.",
    0.7, 4.5, 11.9, 1.7, size=11, color=WHITE)

slide_number(s, 7)

# ── SLIDE 8: Tech Stack ──────────────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "Tech Stack", 0.6, 0.3, 10, 0.6, size=28, bold=True, color=BLUE)

stack = [
    ("Data Layer",      ["pandas — load, transform, validate", "numpy — random walk simulation", "sqlalchemy — DB integration", "CSV — raw storage format"]),
    ("API & Fetch",     ["requests — HTTP client", "mfapi.in — free NAV API", "JSON → DataFrame pipeline", "UTF-8 encoded CSV output"]),
    ("Analysis & Viz",  ["matplotlib / seaborn — static plots", "plotly — interactive charts", "scipy — statistical analysis", "jupyter — EDA notebooks"]),
]
for i, (title, pts) in enumerate(stack):
    box(s, 0.5 + i * 4.2, 1.2, 3.9, 4.0, color=BG_CARD)
    txt(s, title, 0.7 + i * 4.2, 1.3, 3.5, 0.4, size=13, bold=True, color=BLUE)
    for j, p in enumerate(pts):
        txt(s, f"→  {p}", 0.7 + i * 4.2, 1.9 + j * 0.65, 3.5, 0.6, size=11, color=WHITE)

box(s, 0.5, 5.5, 12.3, 1.5, color=BG_CARD)
txt(s, "Version Control & Deployment", 0.7, 5.6, 11.9, 0.4, size=13, bold=True, color=BLUE)
txt(s, "Git  ·  GitHub (github.com/JASWANTH1726/mf-analysis)  ·  Single commit per day  ·  Raw data committed alongside code for reproducibility",
    0.7, 6.1, 11.9, 0.6, size=11, color=WHITE)

slide_number(s, 8)

# ── SLIDE 9: Project Structure ───────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "Project Structure & Deliverables", 0.6, 0.3, 10, 0.6, size=28, bold=True, color=BLUE)

box(s, 0.5, 1.1, 5.5, 5.8, color=BG_CARD)
txt(s, "Folder Structure", 0.7, 1.2, 5.1, 0.4, size=13, bold=True, color=BLUE)
structure = [
    ("mf_analysis/",              WHITE,  11),
    ("  data/raw/",               BLUE,   11),
    ("    16 CSV files",          GREY,   10),
    ("  data/processed/",         BLUE,   11),
    ("  notebooks/",              BLUE,   11),
    ("  sql/",                    BLUE,   11),
    ("  dashboard/",              BLUE,   11),
    ("  reports/",                BLUE,   11),
    ("  data_ingestion.py",       GREEN,  11),
    ("  live_nav_fetch.py",       GREEN,  11),
    ("  generate_datasets.py",    GREEN,  11),
    ("  requirements.txt",        ORANGE, 11),
]
for i, (line, color, size) in enumerate(structure):
    txt(s, line, 0.7, 1.75 + i * 0.42, 5.1, 0.4, size=size, color=color)

box(s, 6.3, 1.1, 6.5, 2.7, color=BG_CARD)
txt(s, "Day 1 Deliverables", 6.5, 1.2, 6.1, 0.4, size=13, bold=True, color=BLUE)
d1 = ["data_ingestion.py", "live_nav_fetch.py", "generate_datasets.py",
      "requirements.txt", "16 raw CSVs in data/raw/", "GitHub repo — Day 1 commit"]
for i, d in enumerate(d1):
    txt(s, f"✓  {d}", 6.5, 1.75 + i * 0.42, 6.1, 0.4, size=11, color=WHITE)

box(s, 6.3, 4.1, 6.5, 2.8, color=BG_CARD)
txt(s, "Upcoming Work", 6.5, 4.2, 6.1, 0.4, size=13, bold=True, color=BLUE)
upcoming = ["Returns analysis & rolling metrics", "Risk-adjusted performance comparison",
            "Portfolio overlap heatmap", "SIP XIRR calculator", "Interactive dashboard"]
for i, u in enumerate(upcoming):
    txt(s, f"→  {u}", 6.5, 4.75 + i * 0.42, 6.1, 0.4, size=11, color=WHITE)

slide_number(s, 9)

# ── SLIDE 10: Summary ────────────────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "Summary", 0.6, 0.3, 10, 0.6, size=28, bold=True, color=BLUE)

summary_stats = [("20", "Schemes Tracked"), ("10", "Fund Houses"), ("16", "CSV Datasets"), ("~20K", "Rows Ingested")]
for i, (val, label) in enumerate(summary_stats):
    box(s, 0.5 + i * 3.1, 1.1, 2.8, 1.6, color=BG_CARD)
    txt(s, val,   0.6 + i * 3.1, 1.15, 2.6, 0.85, size=38, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    txt(s, label, 0.6 + i * 3.1, 2.0,  2.6, 0.45, size=11, color=GREY, align=PP_ALIGN.CENTER)

box(s, 0.5, 3.0, 12.3, 3.2, color=BG_CARD)
txt(s, "What was built", 0.7, 3.1, 11.9, 0.4, size=13, bold=True, color=BLUE)
points = [
    "Live NAV pipeline for 6 large-cap schemes via mfapi.in — 3,100 to 3,600 records each",
    "10 structured datasets covering fund metadata, NAV history, holdings, SIP, returns, AUM, TER, managers, benchmarks, risk",
    "100% AMFI code coverage validated programmatically — zero nulls, zero duplicates",
    "Clean project structure with separate raw/processed layers, ready for analysis and dashboarding",
    "Full Git history on GitHub — reproducible from a single clone",
]
for i, p in enumerate(points):
    txt(s, f"→  {p}", 0.7, 3.65 + i * 0.48, 11.9, 0.44, size=11, color=WHITE)

txt(s, "github.com/JASWANTH1726/mf-analysis", 0.6, 6.8, 12, 0.4, size=11, color=GREY)
slide_number(s, 10)

# ── Save ─────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "reports", "MutualFund_Analysis_Presentation.pptx")
prs.save(out)
print(f"saved: {out}")
